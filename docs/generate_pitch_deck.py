import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_pitch_deck(output_path="FairChain_Pitch_Deck.pptx"):
    prs = Presentation()
    
    # Define slide layouts
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]

    # Slide 1: Title
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "FairChain"
    subtitle.text = "Fair & Predictive Supply Chain Intelligence\n\nGarud (Data Engineer & LLM Integration)\nNirmayee (ML & Dashboarding)"

    # Slide 2: Problem - Supply Chain Unreliability
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "The Massive Cost of Supply Chain Unreliability"
    tf = body_shape.text_frame
    tf.text = "India faces massive logistical bottlenecks due to unpredicted disruptions:"
    p = tf.add_paragraph()
    p.text = "Floods and extreme weather events paralyze critical corridors."
    p = tf.add_paragraph()
    p.text = "Reactive measures lead to huge inventory holding costs and spoilage."
    p = tf.add_paragraph()
    p.text = "Billions of dollars lost annually to supply chain friction."

    # Slide 3: The Hidden Bias
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "The Hidden Bias in Logistics AI"
    tf = body_shape.text_frame
    tf.text = "Current routing software isn't just slow; it's unfair."
    p = tf.add_paragraph()
    p.text = "Models favor large Enterprise Transporters due to higher historical data volume."
    p = tf.add_paragraph()
    p.text = "Local SME Transporters are unfairly penalized by 'low confidence' scores."
    p = tf.add_paragraph()
    p.text = "Result: SMEs are locked out of emergency logistics contracts despite identical capabilities."

    # Slide 4: The FairChain Advantage
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "The FairChain Advantage"
    tf = body_shape.text_frame
    tf.text = "The FairChain Advantage: We combine real-time anomaly detection with fairness-aware routing to predict disruptions before they happen, while ensuring equitable distribution of logistics contracts to SMEs."

    # Slide 5: Our 4-Layer Technical Stack
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Architecture: 4-Layer Stack"
    tf = body_shape.text_frame
    tf.text = "1. Data Layer: Telematics, IMD Weather APIs, OSM routing data."
    p = tf.add_paragraph()
    p.text = "2. ML Engine: Isolation Forests for Anomaly Detection + Fairness Scorecards."
    p = tf.add_paragraph()
    p.text = "3. Explanation Broker: Google Gemini API turning ML scores into human-readable impact."
    p = tf.add_paragraph()
    p.text = "4. Presentation: Mapbox & Next.js Geospatial Dashboard."

    # Slide 6: Fairness Scores
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Fairness Scorecards"
    tf = body_shape.text_frame
    tf.text = "We actively debias AI routing decisions:"
    p = tf.add_paragraph()
    p.text = "Normalize ML confidence scores based on vendor size."
    p = tf.add_paragraph()
    p.text = "Provide transparent 'Fairness Scores' for each recommendation."
    p = tf.add_paragraph()
    p.text = "Prevent monopoly by enterprise carriers during crisis rerouting."

    # Slide 7: Quantified Efficiency Gain
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Quantified Efficiency & Impact"
    tf = body_shape.text_frame
    tf.text = "T-4 Hour Advantage: Our models alert managers 4 hours before official road closures."
    p = tf.add_paragraph()
    p.text = "22% Cost Reduction during emergency rerouting by using qualified local SMEs."
    p = tf.add_paragraph()
    p.text = "Open Innovation: An AI system designed from the ground up for equitable growth."

    # Slide 8: Future Roadmap
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Future Roadmap"
    tf = body_shape.text_frame
    tf.text = "Integration with Port Authorities for real-time maritime capacity."
    p = tf.add_paragraph()
    p.text = "Partnerships with Insurance Firms for dynamic risk premiums."
    p = tf.add_paragraph()
    p.text = "Expansion to Multi-Modal Transport (Rail + Road)."

    prs.save(output_path)
    print(f"Presentation saved successfully to {output_path}")

if __name__ == '__main__':
    create_pitch_deck('FairChain_Pitch_Deck.pptx')
